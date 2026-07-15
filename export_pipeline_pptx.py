"""
PharmGraphRAG Pipeline Diagram — Editable PPTX
Run: python export_pipeline_pptx.py
Output: charts/fig8_pipeline_diagram.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from lxml import etree

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), 'charts')
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ── Helpers ──────────────────────────────────────────────────────────────────
def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def inches(n):
    return Inches(n)

def set_shape_fill(shape, fill_hex, edge_hex, edge_pt=1.8):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.color.rgb = rgb(edge_hex)
    shape.line.width = Pt(edge_pt)

def add_para(tf, text, bold=False, italic=False, size_pt=10,
             color_hex='#374151', align=PP_ALIGN.LEFT, space_before=0):
    para = tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(space_before)
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size_pt)
    run.font.color.rgb = rgb(color_hex)
    return para

def add_connector_with_arrow(slide, x1, y1, x2, y2,
                              color_hex='#6b7280', width_pt=2.0):
    """Add a straight connector and attach an arrowhead at the end via XML."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        inches(x1), inches(y1), inches(x2), inches(y2)
    )
    conn.line.color.rgb = rgb(color_hex)
    conn.line.width = Pt(width_pt)

    # Patch arrowhead at tail end (destination) via OOXML
    sp = conn._element
    spPr = sp.find(qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    # Remove existing end elements to avoid duplicates
    for tag in [qn('a:headEnd'), qn('a:tailEnd')]:
        for el in ln.findall(tag):
            ln.remove(el)
    head = etree.SubElement(ln, qn('a:headEnd'))
    head.set('type', 'none')
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'arrow')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return conn

def add_label(slide, x, y, w, h, text,
              size_pt=8, color_hex='#9ca3af', italic=True):
    txb = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)


# ════════════════════════════════════════════════════════════════════════════
#  Build slide
# ════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)   # widescreen 16:9
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank
slide = prs.slides.add_slide(blank_layout)

# ── Colors ────────────────────────────────────────────────────────────────
C_FILL_LIGHT  = '#dbeafe'
C_FILL_S2     = '#bfdbfe'
C_EDGE        = '#93c5fd'
C_TITLE_DARK  = '#1e3a5f'
C_TEXT        = '#374151'
C_BULLET      = '#4b5563'
C_ARROW       = '#6b7280'
C_AGENTIC_BG  = '#f0f9ff'
C_AGENTIC_EDG = '#7dd3fc'

# ════════════════════════════════════════════════════════════════════════════
#  Title
# ════════════════════════════════════════════════════════════════════════════
title_box = slide.shapes.add_textbox(
    inches(0.3), inches(0.12), inches(12.73), inches(0.55)
)
tf = title_box.text_frame
para = tf.paragraphs[0]
para.alignment = PP_ALIGN.CENTER
run = para.add_run()
run.text = 'Core Algorithm Engine: PharmGraphRAG Pipeline'
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = rgb(C_TITLE_DARK)

# ════════════════════════════════════════════════════════════════════════════
#  S1 Box — top center
# ════════════════════════════════════════════════════════════════════════════
s1 = slide.shapes.add_shape(
    1,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE = 5, but use 1 for rectangle;
        # we'll use freeform rounded rect via autoshape type 5
    inches(2.4), inches(0.78), inches(8.53), inches(2.35)
)
# Use autoshape rounded rect
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches as In

# Re-add as rounded rect (auto shape index 5)
slide.shapes._spTree.remove(s1._element)

def add_rounded_rect(slide, x, y, w, h, fill_hex, edge_hex, edge_pt=1.8, radius_adj=0):
    """Add a rounded rectangle auto-shape."""
    from pptx.oxml.ns import nsmap
    sp = slide.shapes.add_shape(5, inches(x), inches(y), inches(w), inches(h))
    set_shape_fill(sp, fill_hex, edge_hex, edge_pt)
    # Adjust corner radius (0=square, 50000=max round)
    # prstGeom has an adj value
    spPr = sp._element.find(qn('p:spPr'))
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        # Remove existing gd elements
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {16000 + radius_adj}')
    sp.text_frame.word_wrap = True
    return sp

# ── S1 ──
s1 = add_rounded_rect(slide, 2.4, 0.78, 8.53, 2.35, C_FILL_LIGHT, C_EDGE, edge_pt=1.8)
tf1 = s1.text_frame
tf1.margin_left   = Pt(14)
tf1.margin_top    = Pt(10)
tf1.margin_right  = Pt(10)
tf1.margin_bottom = Pt(6)

# First para already exists
p0 = tf1.paragraphs[0]
p0.alignment = PP_ALIGN.LEFT
r0 = p0.add_run()
r0.text = 'S1:'
r0.font.size = Pt(13)
r0.font.bold = True
r0.font.color.rgb = rgb(C_TITLE_DARK)

add_para(tf1, 'Bottom-Up Semantic Retrieval', bold=True, size_pt=15,
         color_hex=C_TITLE_DARK, space_before=2)
add_para(tf1, '', size_pt=5)
add_para(tf1, '\u2022  Always-on hybrid recall (FAISS + entity vector + text match)',
         size_pt=10, color_hex=C_BULLET)
add_para(tf1, '\u2022  Foundational semantic accuracy layer',
         size_pt=10, color_hex=C_BULLET)
add_para(tf1, '\u2022  Vector similarity search + HyDE expansion',
         size_pt=10, color_hex=C_BULLET)
add_para(tf1, '\u2022  Sibling chunk auto-expansion for short clauses',
         size_pt=10, color_hex=C_BULLET)

# ── Fan-out arrows from S1 ──
# S1 bottom center → branch point
S1_CX = 2.4 + 8.53/2          # 6.665"
S1_BOT = 0.78 + 2.35          # 3.13"
BRANCH_Y = 3.52
S2_CX  = 0.2 + 5.9/2          # 3.15"
S3_CX  = 7.23 + 5.9/2         # 10.18"
S2_TOP = 3.68
S3_TOP = 3.68

# Vertical stem S1 → branch
add_connector_with_arrow(slide, S1_CX, S1_BOT, S1_CX, BRANCH_Y - 0.05, C_ARROW, 2.0)

# Horizontal branch line (plain, no arrow)
branch_conn = slide.shapes.add_connector(
    MSO_CONNECTOR_TYPE.STRAIGHT,
    inches(S2_CX), inches(BRANCH_Y),
    inches(S3_CX), inches(BRANCH_Y)
)
branch_conn.line.color.rgb = rgb(C_ARROW)
branch_conn.line.width = Pt(2.0)

# Down arrows to S2 and S3
add_connector_with_arrow(slide, S2_CX, BRANCH_Y, S2_CX, S2_TOP, C_ARROW, 2.0)
add_connector_with_arrow(slide, S3_CX, BRANCH_Y, S3_CX, S3_TOP, C_ARROW, 2.0)

# Labels on arrows
add_label(slide, S1_CX - 1.5, S1_BOT + 0.02, 1.3, 0.35,
          'always executed', size_pt=8, color_hex='#9ca3af')
add_label(slide, (S2_CX + S1_CX)/2 - 0.85, BRANCH_Y - 0.3, 1.4, 0.28,
          'parallel', size_pt=8, color_hex='#9ca3af')
add_label(slide, (S3_CX + S1_CX)/2 - 0.1, BRANCH_Y - 0.3, 1.4, 0.28,
          'parallel', size_pt=8, color_hex='#9ca3af')

# ════════════════════════════════════════════════════════════════════════════
#  S2 Box — bottom left
# ════════════════════════════════════════════════════════════════════════════
s2 = add_rounded_rect(slide, 0.2, 3.68, 5.9, 2.75, C_FILL_S2, C_EDGE, edge_pt=2.2)
tf2 = s2.text_frame
tf2.margin_left   = Pt(14)
tf2.margin_top    = Pt(10)
tf2.margin_right  = Pt(10)
tf2.margin_bottom = Pt(6)

p0 = tf2.paragraphs[0]
p0.alignment = PP_ALIGN.LEFT
r0 = p0.add_run()
r0.text = 'S2:'
r0.font.size = Pt(13)
r0.font.bold = True
r0.font.color.rgb = rgb(C_TITLE_DARK)

add_para(tf2, 'Top-Down Structural Navigation', bold=True, size_pt=15,
         color_hex=C_TITLE_DARK, space_before=2)
add_para(tf2, '', size_pt=5)
add_para(tf2, '\u2022  Dynamic strategy routing', size_pt=10, color_hex=C_BULLET)
add_para(tf2, '\u2022  Intelligent router mechanism', size_pt=10, color_hex=C_BULLET)
add_para(tf2, '\u2022  Graph topology exploration', size_pt=10, color_hex=C_BULLET)
add_para(tf2, '\u2022  Chapter-level document reading', size_pt=10, color_hex=C_BULLET)

# ════════════════════════════════════════════════════════════════════════════
#  S3 Box — bottom right
# ════════════════════════════════════════════════════════════════════════════
s3 = add_rounded_rect(slide, 7.23, 3.68, 5.9, 2.75, C_FILL_LIGHT, C_EDGE, edge_pt=1.8)
tf3 = s3.text_frame
tf3.margin_left   = Pt(14)
tf3.margin_top    = Pt(10)
tf3.margin_right  = Pt(10)
tf3.margin_bottom = Pt(6)

p0 = tf3.paragraphs[0]
p0.alignment = PP_ALIGN.LEFT
r0 = p0.add_run()
r0.text = 'S3:'
r0.font.size = Pt(13)
r0.font.bold = True
r0.font.color.rgb = rgb(C_TITLE_DARK)

add_para(tf3, 'LLM-Guided Graph Walk', bold=True, size_pt=15,
         color_hex=C_TITLE_DARK, space_before=2)
add_para(tf3, '', size_pt=5)
add_para(tf3, '\u2022  Conditionally triggered', size_pt=10, color_hex=C_BULLET)
add_para(tf3, '     (high-risk intent detected)', size_pt=9.5,
         color_hex='#9ca3af', italic=True)
add_para(tf3, '\u2022  Context-sensitive path selection', size_pt=10, color_hex=C_BULLET)
add_para(tf3, '\u2022  Replaces traditional MCTS', size_pt=10, color_hex=C_BULLET)

# ════════════════════════════════════════════════════════════════════════════
#  Merge arrows S2+S3 → Agentic bar
# ════════════════════════════════════════════════════════════════════════════
S2_BOT = 3.68 + 2.75   # 6.43"
S3_BOT = 3.68 + 2.75
MERGE_Y = 6.56
AGENTIC_TOP = 6.62

add_connector_with_arrow(slide, S2_CX, S2_BOT, S2_CX, MERGE_Y - 0.03, C_ARROW, 2.0)
add_connector_with_arrow(slide, S3_CX, S3_BOT, S3_CX, MERGE_Y - 0.03, C_ARROW, 2.0)

horiz = slide.shapes.add_connector(
    MSO_CONNECTOR_TYPE.STRAIGHT,
    inches(S2_CX), inches(MERGE_Y),
    inches(S3_CX), inches(MERGE_Y)
)
horiz.line.color.rgb = rgb(C_ARROW)
horiz.line.width = Pt(2.0)

add_connector_with_arrow(slide, S1_CX, MERGE_Y, S1_CX, AGENTIC_TOP, C_ARROW, 2.0)

# ════════════════════════════════════════════════════════════════════════════
#  Agentic Execution bar — bottom
# ════════════════════════════════════════════════════════════════════════════
ag = add_rounded_rect(slide, 0.2, 6.62, 12.93, 0.78, C_AGENTIC_BG, C_AGENTIC_EDG, edge_pt=2.0)
tf_ag = ag.text_frame
tf_ag.margin_left   = Pt(16)
tf_ag.margin_top    = Pt(8)
tf_ag.margin_right  = Pt(10)
tf_ag.margin_bottom = Pt(4)

p0 = tf_ag.paragraphs[0]
p0.alignment = PP_ALIGN.LEFT
r_bold = p0.add_run()
r_bold.text = 'Agentic Execution    '
r_bold.font.size = Pt(13.5)
r_bold.font.bold = True
r_bold.font.color.rgb = rgb(C_TITLE_DARK)
r_body = p0.add_run()
r_body.text = ('Manufacturer, Regulator, Distributor agents query RAG interface'
               '  \u2192  Unscripted GMP-compliant decisions')
r_body.font.size = Pt(10.5)
r_body.font.color.rgb = rgb(C_BULLET)

# ════════════════════════════════════════════════════════════════════════════
#  Save
# ════════════════════════════════════════════════════════════════════════════
out = os.path.join(OUTPUT_DIR, 'fig8_pipeline_diagram.pptx')
prs.save(out)
print(f'[OK] Saved: {out}')
print('Open in PowerPoint — all text boxes and shapes are fully editable.')
